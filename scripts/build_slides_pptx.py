#!/usr/bin/env python
"""Build the May 26 final-presentation PPTX from the SLIDES.md draft.

Run:
    uv run --with python-pptx python scripts/build_slides_pptx.py

Output: deliverables/final-presentation/Final Presentation.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUTPUT = (
    Path(__file__).resolve().parents[1] / "deliverables/final-presentation/Final Presentation.pptx"
)

SLATE = RGBColor(0x2D, 0x2D, 0x32)
SLATE_LIGHT = RGBColor(0x3A, 0x3A, 0x40)
YELLOW = RGBColor(0xF5, 0xC8, 0x4C)
ORANGE = RGBColor(0xF0, 0x8A, 0x4B)
RED = RGBColor(0xE0, 0x60, 0x6B)
GREEN = RGBColor(0x6E, 0xC4, 0x7A)
CREAM = RGBColor(0xF2, 0xEE, 0xE5)
DIM = RGBColor(0xB8, 0xB8, 0xB8)
BLUE = RGBColor(0x6A, 0xA8, 0xE0)

FOOTER_TEXT = "CSEN 346 · Natural Language Processing · Cyberdyne · May 26, 2026"


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_strip(slide, left, top, width, height, color):
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    strip.fill.solid()
    strip.fill.fore_color.rgb = color
    strip.line.fill.background()
    return strip


def add_text(
    slide,
    text,
    *,
    left,
    top,
    width,
    height,
    size=18,
    bold=False,
    italic=False,
    color=CREAM,
    align=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_bullets(
    slide,
    items,
    *,
    left,
    top,
    width,
    height,
    size=16,
    color=CREAM,
    bullet_color=YELLOW,
    line_spacing=1.15,
):
    """Items: list of strings (plain) or dicts {text, bold, color, indent}."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            spec = {"text": item}
        else:
            spec = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        bullet_run = p.add_run()
        bullet_run.text = spec.get("bullet", "• ")
        bullet_run.font.name = "Calibri"
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = spec.get("bullet_color", bullet_color)
        bullet_run.font.bold = True
        run = p.add_run()
        run.text = spec["text"]
        run.font.name = "Calibri"
        run.font.size = Pt(spec.get("size", size))
        run.font.bold = spec.get("bold", False)
        run.font.italic = spec.get("italic", False)
        run.font.color.rgb = spec.get("color", color)
    return box


def add_footer(slide):
    add_text(
        slide,
        FOOTER_TEXT,
        left=Inches(0.4),
        top=Inches(7.05),
        width=Inches(12.0),
        height=Inches(0.3),
        size=10,
        color=DIM,
    )


def add_title(slide, text, *, size=32, color=YELLOW, top=Inches(0.35)):
    add_text(
        slide,
        text,
        left=Inches(0.5),
        top=top,
        width=Inches(12.0),
        height=Inches(0.8),
        size=size,
        bold=True,
        color=color,
    )
    add_accent_strip(slide, Inches(0.5), Inches(1.05), Inches(1.2), Inches(0.05), YELLOW)


def add_speaker_label(slide, speaker, time_str, *, top=Inches(1.18)):
    add_text(
        slide,
        f"{speaker} · {time_str}",
        left=Inches(0.5),
        top=top,
        width=Inches(8.0),
        height=Inches(0.3),
        size=12,
        color=DIM,
        font_name="Calibri",
    )


def set_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def add_table(
    slide,
    headers,
    rows,
    *,
    left,
    top,
    width,
    height,
    header_color=YELLOW,
    header_fill=SLATE_LIGHT,
    body_color=CREAM,
    row_fills=None,
    font_size=12,
):
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = tbl_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.clear()
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = header_color
        run.font.name = "Calibri"
    for i, row in enumerate(rows, start=1):
        fill = row_fills[i - 1] if row_fills and i - 1 < len(row_fills) else SLATE
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            tf = cell.text_frame
            tf.clear()
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            if isinstance(val, dict):
                run = p.add_run()
                run.text = val["text"]
                run.font.size = Pt(val.get("size", font_size))
                run.font.bold = val.get("bold", False)
                run.font.italic = val.get("italic", False)
                run.font.color.rgb = val.get("color", body_color)
                run.font.name = "Calibri"
            else:
                run = p.add_run()
                run.text = str(val)
                run.font.size = Pt(font_size)
                run.font.color.rgb = body_color
                run.font.name = "Calibri"
    return tbl_shape


def make_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---------- SLIDE 1 — Title ----------
    s = prs.slides.add_slide(blank)
    set_background(s, SLATE)
    add_accent_strip(s, 0, Inches(0.0), Inches(13.333), Inches(0.12), YELLOW)
    add_accent_strip(s, 0, Inches(0.12), Inches(13.333), Inches(0.04), ORANGE)
    add_accent_strip(s, 0, Inches(0.16), Inches(13.333), Inches(0.04), RED)
    add_text(
        s,
        "Beyond Memorization",
        left=Inches(0.6),
        top=Inches(1.6),
        width=Inches(12.0),
        height=Inches(1.0),
        size=54,
        bold=True,
        color=CREAM,
    )
    add_text(
        s,
        "Reproducing and Extending KELE for Socratic Teaching with LLMs",
        left=Inches(0.6),
        top=Inches(2.5),
        width=Inches(12.0),
        height=Inches(0.6),
        size=24,
        color=YELLOW,
    )
    add_accent_strip(s, Inches(0.6), Inches(3.2), Inches(1.5), Inches(0.05), YELLOW)
    add_text(
        s,
        "Team Cyberdyne — Ulises Chavarria · Maximilian Khan",
        left=Inches(0.6),
        top=Inches(3.4),
        width=Inches(12.0),
        height=Inches(0.4),
        size=20,
        color=CREAM,
    )
    add_text(
        s,
        "CSEN 346 · Santa Clara University · Spring 2026",
        left=Inches(0.6),
        top=Inches(3.85),
        width=Inches(12.0),
        height=Inches(0.4),
        size=16,
        color=DIM,
    )
    add_text(
        s,
        "Paper reproduced: Peng et al., KELE: A Multi-Agent Framework for Structured Socratic Teaching, EMNLP 2025 Findings",
        left=Inches(0.6),
        top=Inches(5.2),
        width=Inches(12.0),
        height=Inches(0.4),
        size=14,
        color=DIM,
    )
    add_text(
        s,
        "github.com/ulises-c/csen-346  ·  huggingface.co/ulises-c",
        left=Inches(0.6),
        top=Inches(5.65),
        width=Inches(12.0),
        height=Inches(0.4),
        size=14,
        color=BLUE,
    )
    set_notes(
        s,
        "Hi everyone — we're team Cyberdyne. I'm Ulises, this is Max. "
        "We spent the quarter reproducing a 2025 EMNLP Findings paper called KELE — "
        "a multi-agent framework for Socratic teaching with LLMs — and then extending "
        "it in three ways. The most surprising thing we found is that the paper's "
        "headline result reverses under a fair metric. We'll walk you through what we "
        "built, what we measured, and what that reversal actually means.\n\n"
        'Transition: "Let me start with a quick refresher on the problem and what KELE proposes."',
    )

    # ---------- SLIDE 2 — Refresher ----------
    s = prs.slides.add_slide(blank)
    set_background(s, SLATE)
    add_title(s, "Quick Refresher: KELE in One Slide")
    add_speaker_label(s, "Ulises", "~60 s · you saw this in our April pitch")
    add_footer(s)

    add_text(
        s,
        "The Socratic problem: guide the student to the answer; don't give it.",
        left=Inches(0.5),
        top=Inches(1.7),
        width=Inches(12.0),
        height=Inches(0.4),
        size=18,
        color=CREAM,
    )

    add_text(
        s,
        "SocRules — KELE's 5-stage framework:",
        left=Inches(0.5),
        top=Inches(2.3),
        width=Inches(6.0),
        height=Inches(0.4),
        size=16,
        bold=True,
        color=YELLOW,
    )
    # 5-stage strip
    stages = [
        ("a", "greeting"),
        ("b", "questioning"),
        ("c", "correction"),
        ("d", "affirmation"),
        ("e", "closure"),
    ]
    box_w = Inches(2.3)
    gap = Inches(0.1)
    start_x = Inches(0.5)
    for i, (letter, name) in enumerate(stages):
        x = start_x + (box_w + gap) * i
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.8), box_w, Inches(0.7))
        b.fill.solid()
        b.fill.fore_color.rgb = SLATE_LIGHT
        b.line.color.rgb = YELLOW
        b.line.width = Pt(1)
        tf = b.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{letter})  {name}"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = CREAM
        run.font.name = "Calibri"

    add_text(
        s,
        "KELE architecture (Peng et al., EMNLP 2025 Findings):",
        left=Inches(0.5),
        top=Inches(3.85),
        width=Inches(12.0),
        height=Inches(0.4),
        size=16,
        bold=True,
        color=YELLOW,
    )
    add_bullets(
        s,
        [
            {
                "text": "Consultant (GPT-4o) — classifies which SocRule stage the dialogue is in",
                "size": 15,
            },
            {
                "text": "Teacher (SocratTeachLLM, GLM-4 9B fine-tuned) — generates the response",
                "size": 15,
            },
            {
                "text": "Benchmark: SocratDataset — 6,803 Chinese dialogues · 681-dialogue test split",
                "size": 15,
            },
        ],
        left=Inches(0.5),
        top=Inches(4.3),
        width=Inches(12.0),
        height=Inches(1.4),
    )
    add_text(
        s,
        'Published headline: "SocratTeachLLM surpasses GPT-4o on ROUGE-1/2 and BLEU-4."',
        left=Inches(0.5),
        top=Inches(5.85),
        width=Inches(12.0),
        height=Inches(0.5),
        size=16,
        bold=True,
        italic=True,
        color=ORANGE,
    )
    set_notes(
        s,
        "Quick refresher — most of you saw our initial pitch in April, so I'll move fast. "
        "The Socratic teaching problem is to guide the student to the answer instead of "
        "giving it. KELE operationalizes that with a fixed five-stage script — greeting, "
        "questioning, correction, affirmation, closure — which they call SocRules.\n\n"
        "The architecture is two agents: a Consultant that classifies which stage the "
        "dialogue is in — GPT-4o in the paper — and a Teacher that generates the response "
        "— SocratTeachLLM, a GLM-4 9B fine-tuned on their dataset. The benchmark is "
        "SocratDataset, about 6,800 Chinese tutoring dialogues with a 681-dialogue test "
        "split. The published headline is that their 9B fine-tuned Teacher surpasses "
        "GPT-4o on ROUGE-1, ROUGE-2, and BLEU-4. Note I said 'on their metrics' — that's "
        "the thread I'm pulling in two slides.\n\n"
        'Transition: "In April we proposed four future improvements. Here\'s the scoreboard."',
    )

    # ---------- SLIDE 3 — Pitch -> Delivery Scoreboard ----------
    s = prs.slides.add_slide(blank)
    set_background(s, SLATE)
    add_title(s, "Pitch → Delivery Scoreboard")
    add_speaker_label(s, "Ulises", "~90 s")
    add_footer(s)

    add_text(
        s,
        "In April we proposed four future improvements. Here's what actually happened:",
        left=Inches(0.5),
        top=Inches(1.55),
        width=Inches(12.5),
        height=Inches(0.4),
        size=14,
        italic=True,
        color=DIM,
    )

    add_table(
        s,
        headers=["April pitch direction", "Outcome"],
        rows=[
            [
                "1. MoE for multi-subject teaching",
                {"text": "❌  Dropped — no multi-subject corpus available", "color": RED},
            ],
            [
                "2. RL-based stage transitions",
                {
                    "text": "❌  Dropped — bottleneck wasn't transitions; it was the consultant",
                    "color": RED,
                },
            ],
            [
                {"text": "3. Learned consultant (discrete classifier replacing LLM)", "bold": True},
                {
                    "text": "✅  Shipped — became our locked headline",
                    "color": GREEN,
                    "bold": True,
                },
            ],
            [
                {"text": "4. Stronger evaluation (ablation + cross-domain)", "bold": True},
                {
                    "text": "✅  Shipped — significantly expanded",
                    "color": GREEN,
                    "bold": True,
                },
            ],
        ],
        left=Inches(0.5),
        top=Inches(2.0),
        width=Inches(12.3),
        height=Inches(2.1),
        font_size=14,
    )

    add_text(
        s,
        "Plus three contributions we didn't see coming:",
        left=Inches(0.5),
        top=Inches(4.4),
        width=Inches(12.0),
        height=Inches(0.4),
        size=16,
        bold=True,
        color=YELLOW,
    )
    # 3 cards in a row
    card_titles = [
        ("🔍  Methodological critique", "ROUGE/BLEU invert the true ranking (slide 4)"),
        ("📐  Unified score", "memorization-resistant ranking (slide 5)"),
        ("🌐  SocratDataset-EN", "full English translation + cross-lingual transfer (slide 7)"),
    ]
    card_w = Inches(4.0)
    card_gap = Inches(0.15)
    card_start = Inches(0.5)
    for i, (title, body) in enumerate(card_titles):
        x = card_start + (card_w + card_gap) * i
        card = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.95), card_w, Inches(1.55)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = SLATE_LIGHT
        card.line.color.rgb = ORANGE
        card.line.width = Pt(1)
        tf = card.text_frame
        tf.clear()
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.12)
        tf.margin_bottom = Inches(0.12)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.bold = True
        r.font.size = Pt(15)
        r.font.color.rgb = YELLOW
        r.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(12)
        r2.font.color.rgb = CREAM
        r2.font.name = "Calibri"

    set_notes(
        s,
        "In April we made four promises about what we'd try. Here's the honest scoreboard.\n\n"
        "We dropped two of them. Mixture-of-Experts for multi-subject teaching — we couldn't "
        "find a multi-subject Socratic corpus within our budget, and without one, MoE has "
        "nothing to gate on. So we dropped it cleanly.\n\n"
        "RL-based stage transitions — we dropped this for a more interesting empirical reason. "
        "Our reproduction work showed that the SocRule transitions are not the bottleneck in "
        "KELE. The bottleneck is the consultant's stage classification. So learning the "
        "transitions wouldn't have moved the needle — improving the consultant did.\n\n"
        "We shipped the other two. Item 3 — the learned consultant via discrete classifier — "
        "became our locked paper headline. Max will show you the numbers. Item 4 — stronger "
        "evaluation — expanded significantly: 142 model configurations measured, an 8-cell "
        "cross-teacher matrix, bootstrap convergence analysis. We over-delivered on "
        "evaluation.\n\n"
        "And on top of those two, we found three things we didn't anticipate when we pitched. "
        "The biggest one — the one I'm going to show you next — is that the paper's published "
        "winner isn't actually the winner under a memorization-resistant metric. That finding "
        "drove everything that came after it.\n\n"
        'Transition: "Let me show you what we found."',
    )

    # ---------- SLIDE 4 — The Inversion ----------
    s = prs.slides.add_slide(blank)
    set_background(s, SLATE)
    add_title(s, "The Benchmark Inversion")
    add_speaker_label(s, "Ulises", "~90 s")
    add_footer(s)

    add_text(
        s,
        "Recall (April pitch): the paper reports STL beats GPT-4o on all 13 metrics — "
        "4 token-overlap + 4 rule-compliance + 5 LLM-judge rubric.",
        left=Inches(0.5),
        top=Inches(1.55),
        width=Inches(12.5),
        height=Inches(0.5),
        size=13,
        italic=True,
        color=DIM,
    )
    add_text(
        s,
        "On the 4 token-overlap metrics, the ranking inverts under a memorization-resistant evaluation.",
        left=Inches(0.5),
        top=Inches(2.0),
        width=Inches(12.5),
        height=Inches(0.4),
        size=15,
        bold=True,
        color=CREAM,
    )

    add_table(
        s,
        headers=["", "Paper's ROUGE/BLEU ranking", "Our unified ranking (pedagogical)"],
        rows=[
            [
                "🥇 #1",
                "GPT-4o + SocratTeachLLM",
                {"text": "Gemma 31B + BERT (ours)", "bold": True, "color": GREEN},
            ],
            ["🥈 #2", "Claude Opus + BERT", "Claude Sonnet + BERT"],
            [
                "🚨 Last",
                "Claude Opus (raw)",
                {"text": "GPT-4o + SocratTeachLLM", "bold": True, "color": RED},
            ],
        ],
        left=Inches(0.5),
        top=Inches(2.55),
        width=Inches(12.3),
        height=Inches(1.6),
        font_size=14,
    )

    add_text(
        s,
        "Why the surface metrics mislead:",
        left=Inches(0.5),
        top=Inches(4.45),
        width=Inches(12.0),
        height=Inches(0.4),
        size=15,
        bold=True,
        color=YELLOW,
    )
    add_bullets(
        s,
        [
            {
                "text": "ROUGE/BLEU reward n-gram overlap with the reference → reward mimicry of training data",
                "size": 13,
            },
            {
                "text": "SocratTeachLLM was fine-tuned on 90% of the dataset (paper §4.3) — and the train/test split was never released",
                "size": 13,
            },
            {
                "text": "Fingerprint: gap to STL widens with n-gram order — ROUGE-1 Δ=+1.59 → ROUGE-2 Δ=+4.92 → BLEU-4 Δ=+4.07. Memorization compounds.",
                "size": 13,
                "color": ORANGE,
            },
            {
                "text": "(The other 9 metrics — slide 7 contamination evidence implicates them too.)",
                "size": 13,
                "color": DIM,
                "italic": True,
            },
        ],
        left=Inches(0.5),
        top=Inches(4.9),
        width=Inches(12.5),
        height=Inches(1.9),
    )
    set_notes(
        s,
        "In April we showed you the paper's table — SocratTeachLLM beats GPT-4o on all 13 "
        "metrics. Four of them are token-overlap — ROUGE-1, ROUGE-2, ROUGE-L, BLEU-4. The "
        "other nine are pedagogical rubric scores. The paper presents this as a clean "
        "9B-beats-frontier story.\n\n"
        "Here's what we found. On the same ten configurations the paper evaluates, those "
        "four token-overlap metrics produce the ranking on the left. SocratTeachLLM first, "
        "Claude Opus raw last. Now look at the right — same ten configurations, ranked "
        "under a memorization-resistant metric. The ordering is almost exactly reversed. "
        "Our open-weight Gemma 31B plus BERT system is first; SocratTeachLLM is last.\n\n"
        "Why does this happen? ROUGE and BLEU are token-overlap metrics. They reward output "
        "that looks lexically similar to the reference. They don't measure whether the "
        "model taught anything. And SocratTeachLLM was fine-tuned on 90% of the dataset — "
        "the paper says so in section 4.3 — but the train/test split was never released. "
        "So any third party drawing a different random 10% sample is measuring memorization, "
        "not teaching.\n\n"
        "The fingerprint is in the numbers. The gap to STL widens as the n-gram order grows. "
        "Unigrams, plus 1.59. Bigrams, plus 4.92. 4-grams, plus 4.07. Memorization compounds "
        "with sequence length. We don't see that on the LLM-judge rubric or on our "
        "memorization-resistant metrics — and Max will show you on slide 7 that when we "
        "evaluate STL on data it has never seen, its score collapses 30 points. That "
        "implicates the broader 13-metric claim, not just ROUGE and BLEU.\n\n"
        'Transition: "So we needed a different metric. Max will explain what we built — '
        'and the new architecture that goes with it."',
    )

    # ---------- SLIDE 5 — BERT Consultant + Unified Metric ----------
    s = prs.slides.add_slide(blank)
    set_background(s, SLATE)
    add_title(s, "BERT Consultant + Unified Metric")
    add_speaker_label(
        s,
        "Max",
        "~90 s · shipping pitch direction #3 (Learned Consultant)",
    )
    add_footer(s)

    add_text(
        s,
        "1. Consultant swap — GPT-4o → fine-tuned bge-small-zh (24M params)",
        left=Inches(0.5),
        top=Inches(1.55),
        width=Inches(12.5),
        height=Inches(0.4),
        size=17,
        bold=True,
        color=YELLOW,
    )
    add_bullets(
        s,
        [
            {"text": "Frozen base + 5-way + 34-way classification heads", "size": 13},
            {
                "text": "Matches or exceeds GPT-4o on state routing — at $0 per inference",
                "size": 13,
            },
        ],
        left=Inches(0.5),
        top=Inches(2.0),
        width=Inches(6.5),
        height=Inches(0.9),
    )
    add_table(
        s,
        headers=["Consultant", "Teacher", "State Acc", "Δ"],
        rows=[
            ["GPT-4o (original KELE)", "SocratTeachLLM", "25.94%", "baseline"],
            [
                {"text": "BERT (ours)", "bold": True},
                "Gemma 31B",
                "48.15%",
                {"text": "+22.21 pp", "bold": True, "color": GREEN},
            ],
            [
                {"text": "BERT (ours)", "bold": True},
                "Claude Sonnet",
                "49.97%",
                {"text": "+24.03 pp", "bold": True, "color": GREEN},
            ],
        ],
        left=Inches(0.5),
        top=Inches(3.0),
        width=Inches(7.5),
        height=Inches(1.4),
        font_size=12,
    )

    add_text(
        s,
        "2. Unified score — memorization-resistant ranking",
        left=Inches(8.3),
        top=Inches(1.55),
        width=Inches(4.8),
        height=Inches(0.4),
        size=17,
        bold=True,
        color=YELLOW,
    )
    formula_box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.3),
        Inches(2.05),
        Inches(4.8),
        Inches(0.75),
    )
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = SLATE_LIGHT
    formula_box.line.color.rgb = YELLOW
    formula_box.line.width = Pt(1)
    tf = formula_box.text_frame
    tf.clear()
    tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "unified = 0.5 × stage_balanced + 0.5 × (judge × 10)"
    r.font.name = "Consolas"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = CREAM
    add_bullets(
        s,
        [
            {
                "text": "stage_balanced = macro-F1 over 5 SocRule stages (closure is rare; macro-acc hides it)",
                "size": 12,
            },
            {
                "text": "LLM-judge = Claude Sonnet 4.6, 4 rubric axes (Socratic validity · stage advance · age-appropriate · question form)",
                "size": 12,
            },
            {
                "text": "Two orthogonal axes — neither rewards token overlap",
                "size": 12,
                "color": ORANGE,
            },
        ],
        left=Inches(8.3),
        top=Inches(2.95),
        width=Inches(4.8),
        height=Inches(3.5),
    )

    add_text(
        s,
        "Bottom line: a 24M-param classifier replaces the LLM consultant + we have a metric the paper's "
        "fine-tuned model can't game.",
        left=Inches(0.5),
        top=Inches(6.2),
        width=Inches(12.5),
        height=Inches(0.6),
        size=14,
        bold=True,
        italic=True,
        color=ORANGE,
    )
    set_notes(
        s,
        "We made two changes to KELE. The first is structural — and this is direction "
        "number 3 from our April pitch, the learned consultant via discrete classifier. "
        "We replaced GPT-4o with a fine-tuned BERT model — specifically bge-small-zh, "
        "24 million parameters, frozen base with two classification heads. A 5-way head "
        "for the SocRule stage, a 34-way head for finer-grained sub-states. Training "
        "takes under a hundred seconds on a consumer GPU.\n\n"
        "The table shows what happens. The original KELE consultant — GPT-4o — gets the "
        "dialogue's state right 26% of the time. Our BERT consultant, paired with Gemma "
        "31B as the teacher, gets it right 48% of the time. With Claude Sonnet, 50%. "
        "That's a 22 to 24 percentage-point absolute lift, and our consultant runs "
        "locally at zero API cost. The intuition is just that classifying which of five "
        "Socratic stages a dialogue is in is a classification problem, not a generation "
        "problem — you don't need a hundred-billion-parameter model for that.\n\n"
        "The second change is the metric. We propose the unified score — half "
        "stage-balanced accuracy, half an LLM judge. Stage-balanced is macro-F1 over the "
        "five SocRule stages — this matters because the stage distribution is skewed; "
        "questioning is about 40% of turns, closure is about 10%. Macro-accuracy hides "
        "failures on the rare stages, and closure is exactly the stage you can't afford "
        "to fail on. The LLM judge is Claude Sonnet 4.6 scoring four rubric axes — "
        "Socratic validity, stage advancement, age-appropriateness, question form. The "
        "two axes are orthogonal — neither rewards token-level mimicry of the training "
        "set, so neither inflates SocratTeachLLM's score.\n\n"
        "Transition: \"Here's what happens when you actually run the full leaderboard on "
        'this metric."',
    )

    # ---------- SLIDE 6 — Master Leaderboard + Parity ----------
    s = prs.slides.add_slide(blank)
    set_background(s, SLATE)
    add_title(s, "Master Leaderboard + Local–Frontier Parity")
    add_speaker_label(
        s,
        "Max",
        "~90 s · shipping pitch direction #4 (Stronger Evaluation)",
    )
    add_footer(s)

    add_text(
        s,
        "142 configurations measured. 37 with full unified score. Top of the pedagogical ranking:",
        left=Inches(0.5),
        top=Inches(1.55),
        width=Inches(12.5),
        height=Inches(0.4),
        size=14,
        italic=True,
        color=DIM,
    )
    add_table(
        s,
        headers=["Rank", "Config", "n", "Unified", "Stage_bal", "Judge"],
        rows=[
            [
                "🥇 1",
                "bert × Gemma-31B · top3 · n=50",
                "278",
                {"text": "70.08", "bold": True, "color": GREEN},
                "58.48",
                "8.17",
            ],
            [
                "🥈 2",
                "bert × Claude-Sonnet · top3 · n=681",
                "3840",
                "70.06",
                "58.17",
                "8.19",
            ],
            [
                "🥉 3",
                "bert × Claude-Opus · fewshot10 · n=50",
                "271",
                "69.79",
                "58.73",
                "8.08",
            ],
            [
                "🏆 7",
                {"text": "bert × Gemma-31B · fewshot10 · n=681", "bold": True},
                "3834",
                {"text": "68.65", "bold": True, "color": YELLOW},
                "55.42",
                "8.19",
            ],
            [
                "⚠️ 8",
                "qwen3.5 × SocratTeachLLM · n=50",
                "288",
                "68.21",
                "63.40",
                "7.30",
            ],
            [
                "⚠️ last",
                "Claude-Sonnet × SocratTeachLLM · EN · n=50",
                "303",
                {"text": "44.36", "color": RED},
                "22.55",
                "6.62",
            ],
        ],
        left=Inches(0.5),
        top=Inches(2.05),
        width=Inches(12.3),
        height=Inches(2.6),
        font_size=12,
    )

    add_text(
        s,
        "🏆 = locked headline (full n=681)    ⚠️ = SocratTeachLLM, contamination-inflated",
        left=Inches(0.5),
        top=Inches(4.85),
        width=Inches(12.0),
        height=Inches(0.4),
        size=12,
        italic=True,
        color=DIM,
    )

    parity_box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(5.45),
        Inches(12.3),
        Inches(1.4),
    )
    parity_box.fill.solid()
    parity_box.fill.fore_color.rgb = SLATE_LIGHT
    parity_box.line.color.rgb = YELLOW
    parity_box.line.width = Pt(1.5)
    tf = parity_box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "PARITY  "
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = YELLOW
    r2 = p.add_run()
    r2.text = (
        "Best frontier (Claude Sonnet, 70.06) vs. best open-weight (Gemma 31B, 70.08) → gap = 0.02 pts. "
        "Locked headline at n=681 trails by 1.41 — inside sampling noise."
    )
    r2.font.size = Pt(13)
    r2.font.color.rgb = CREAM
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    r3 = p2.add_run()
    r3.text = "Open-weight · 32 GB consumer GPU · $0 per inference · indistinguishable from Claude Sonnet on a fair metric."
    r3.font.size = Pt(13)
    r3.font.italic = True
    r3.font.color.rgb = ORANGE
    set_notes(
        s,
        "We ran 142 configurations end-to-end on the SocratDataset test split. 37 of them "
        "have the full unified score with both stage-balanced and LLM-judge measurements. "
        "Here's the top of the table.\n\n"
        "The headline finding is at the top. The best frontier configuration — a BERT "
        "consultant paired with Claude Sonnet as the teacher, full n=681 — scores unified "
        "70.06. The best open-weight configuration — BERT consultant paired with Gemma 31B "
        "— scores 70.08. That's a gap of two hundredths of a point. Effectively zero. "
        "Indistinguishable.\n\n"
        "The row marked with the trophy — number 7 — is our locked paper headline: BERT "
        "consultant, Gemma 31B teacher, ten-shot prompting, full n=681. It scores 68.65. "
        "That's 1.41 points behind the frontier ceiling — but well inside our measurement "
        "noise at this sample size. Importantly, this configuration runs entirely on a "
        "32 GB consumer GPU at zero per-inference API cost.\n\n"
        "The two rows with warning triangles are SocratTeachLLM-based — "
        "contamination-inflated, exactly the pattern from slide 4. And notice the very "
        "last row — Claude Sonnet teaching with SocratTeachLLM on English — scores 44 "
        "unified, dead last. That's a contamination signature flipping the other way when "
        "the test data leaves the training distribution.\n\n"
        "So the takeaway: open-weight local models on a memorization-resistant benchmark "
        "— they don't trail the frontier. They sit right next to it.\n\n"
        'Transition: "But the contamination claim is doing a lot of work on this slide. '
        'Let me show you the direct evidence."',
    )

    # ---------- SLIDE 7 — Contamination + Cross-Lingual ----------
    s = prs.slides.add_slide(blank)
    set_background(s, SLATE)
    add_title(s, "Contamination Evidence + Cross-Lingual Transfer")
    add_speaker_label(s, "Max", "~75 s")
    add_footer(s)

    add_text(
        s,
        "Direct contamination test: evaluate SocratTeachLLM on a fully synthetic test set — "
        "Claude-generated, never seen during training.",
        left=Inches(0.5),
        top=Inches(1.55),
        width=Inches(12.5),
        height=Inches(0.5),
        size=14,
        color=CREAM,
    )

    # Two-bar visualization for the 63.4 → 32.86 collapse
    chart_left = Inches(0.6)
    chart_top = Inches(2.3)
    chart_h = Inches(2.6)
    # bar 1: 63.4 → height proportional
    full_h = chart_h
    bar_w = Inches(1.5)
    gap_x = Inches(0.4)

    h1 = full_h
    bar1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, chart_left, chart_top, bar_w, h1)
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = BLUE
    bar1.line.fill.background()
    add_text(
        s,
        "63.4",
        left=chart_left,
        top=chart_top - Inches(0.45),
        width=bar_w,
        height=Inches(0.4),
        size=22,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        "Original test set\n(stage_bal)",
        left=chart_left,
        top=chart_top + h1 + Inches(0.05),
        width=bar_w,
        height=Inches(0.6),
        size=11,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )

    # bar 2: 32.86 ~ height proportional
    h2 = Emu(int(full_h.emu * 32.86 / 63.4))
    bar2_x = chart_left + bar_w + gap_x
    bar2_top = chart_top + (full_h - h2)
    bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar2_x, bar2_top, bar_w, h2)
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = RED
    bar2.line.fill.background()
    add_text(
        s,
        "32.86",
        left=bar2_x,
        top=bar2_top - Inches(0.45),
        width=bar_w,
        height=Inches(0.4),
        size=22,
        bold=True,
        color=RED,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        "Synthetic clean test\n(unseen data)",
        left=bar2_x,
        top=chart_top + h1 + Inches(0.05),
        width=bar_w,
        height=Inches(0.6),
        size=11,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )

    # delta callout
    delta_box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        chart_left + bar_w * 2 + gap_x + Inches(0.4),
        chart_top + Inches(0.6),
        Inches(2.6),
        Inches(1.4),
    )
    delta_box.fill.solid()
    delta_box.fill.fore_color.rgb = SLATE_LIGHT
    delta_box.line.color.rgb = RED
    delta_box.line.width = Pt(1.5)
    tf = delta_box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "−30 pts"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RED
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "= memorization"
    r2.font.size = Pt(14)
    r2.font.italic = True
    r2.font.color.rgb = CREAM

    # Cross-lingual + PR #79 on the right
    add_text(
        s,
        "Cross-lingual transfer (SocratDataset-EN):",
        left=Inches(8.0),
        top=Inches(2.2),
        width=Inches(5.0),
        height=Inches(0.4),
        size=16,
        bold=True,
        color=YELLOW,
    )
    add_bullets(
        s,
        [
            {"text": "Translated all 6,803 dialogues to English; published on HF", "size": 12},
            {
                "text": "Our qwen3.5 consultant (Chinese-trained only) → unified 65.11 at n=400 on the English split",
                "size": 12,
            },
            {"text": "SocRule routing is language-invariant", "size": 12, "color": ORANGE},
            {
                "text": '(Directly addresses the paper\'s own §Limitations: "No cross-lingual evaluation.")',
                "size": 11,
                "color": DIM,
                "italic": True,
            },
        ],
        left=Inches(8.0),
        top=Inches(2.65),
        width=Inches(5.0),
        height=Inches(2.2),
    )

    add_text(
        s,
        "PR #79 (2026-05-25): contamination probe now bilingual — SocratDataset-SYNTHETIC-EN (n=75) live on HF.",
        left=Inches(0.5),
        top=Inches(6.3),
        width=Inches(12.5),
        height=Inches(0.5),
        size=12,
        italic=True,
        color=DIM,
    )
    set_notes(
        s,
        "To prove the contamination claim directly, we built a fully synthetic test set — "
        "75 Socratic tutoring dialogues generated by Claude, schema-matched to "
        "SocratDataset but never seen by any model during training. We then evaluated "
        "SocratTeachLLM on it. On the original test set, SocratTeachLLM's stage-balanced "
        "accuracy is 63.4 — that's the number the paper reports. On the synthetic clean "
        "test set, it drops to 32.86. A 30-point collapse. That's not a model that "
        "learned to teach — that's a model that learned to recite the training set.\n\n"
        "The second result on this slide is the cross-lingual transfer experiment. We "
        "translated the entire 6,803-dialogue SocratDataset to English — that's "
        "SocratDataset-EN, live on HuggingFace. Then we took our Chinese-trained "
        "consultant — never saw a word of English during training — and evaluated it on "
        "the English split. It scores unified 65.11. The SocRule routing behavior is "
        "essentially language-invariant. The pedagogical stages are universal.\n\n"
        "One quick update from the last week — we just shipped PR #79, which extends the "
        "synthetic contamination probe to English as well, so we now have "
        "memorization-resistant evaluation data in both languages.\n\n"
        'Transition: "Let me close with what to take away."',
    )

    # ---------- SLIDE 8 — Conclusion ----------
    s = prs.slides.add_slide(blank)
    set_background(s, SLATE)
    add_title(s, "Conclusion + Takeaways")
    add_speaker_label(s, "Max", "~75 s")
    add_footer(s)

    findings = [
        (
            "1",
            "Replace the LLM consultant with a 24M-param BERT classifier.",
            "Same state-routing accuracy. Zero per-run API cost.",
        ),
        (
            "2",
            "ROUGE/BLEU are misleading for Socratic teaching.",
            'The paper\'s "winner" ranks last on pedagogical accuracy. The unified score fixes this.',
        ),
        (
            "3",
            "Open-weight local models reach frontier parity on a memorization-resistant benchmark.",
            "Gemma 31B + BERT + 10-shot ≈ Claude Sonnet within sampling noise.",
        ),
    ]
    card_w = Inches(4.0)
    card_gap = Inches(0.15)
    card_start = Inches(0.5)
    for i, (num, head, sub) in enumerate(findings):
        x = card_start + (card_w + card_gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), card_w, Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = SLATE_LIGHT
        card.line.color.rgb = YELLOW
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.clear()
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.18)
        tf.margin_bottom = Inches(0.18)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = num
        r.font.size = Pt(40)
        r.font.bold = True
        r.font.color.rgb = YELLOW
        r.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = head
        r2.font.size = Pt(14)
        r2.font.bold = True
        r2.font.color.rgb = CREAM
        r2.font.name = "Calibri"
        p3 = tf.add_paragraph()
        p3.space_before = Pt(8)
        r3 = p3.add_run()
        r3.text = sub
        r3.font.size = Pt(12)
        r3.font.color.rgb = DIM
        r3.font.italic = True
        r3.font.name = "Calibri"

    pull = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(4.45),
        Inches(12.3),
        Inches(0.9),
    )
    pull.fill.solid()
    pull.fill.fore_color.rgb = SLATE_LIGHT
    pull.line.color.rgb = ORANGE
    pull.line.width = Pt(1.5)
    tf = pull.text_frame
    tf.clear()
    tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '"If you\'re evaluating educational AI, measure pedagogy — not n-gram overlap."'
    r.font.size = Pt(20)
    r.font.italic = True
    r.font.color.rgb = ORANGE
    r.font.name = "Calibri"

    add_text(
        s,
        "In flight for June 4: Stage 2 SFT pipeline shipped (PR #79). Training in progress on "
        "Gemma 4 31B-IT (pivoted from Qwen3.6-27B — ROCm/Triton 3.6.0 issue on gfx1201).",
        left=Inches(0.5),
        top=Inches(5.55),
        width=Inches(12.3),
        height=Inches(0.6),
        size=13,
        color=DIM,
        italic=True,
    )
    add_text(
        s,
        "github.com/ulises-c/csen-346  ·  huggingface.co/ulises-c  ·  paper draft in deliverables/overleaf/",
        left=Inches(0.5),
        top=Inches(6.4),
        width=Inches(12.3),
        height=Inches(0.4),
        size=14,
        color=BLUE,
    )
    set_notes(
        s,
        "Three things to take home. First: you can replace a hundred-billion-parameter LLM "
        "consultant with a 24-million-parameter BERT classifier and you don't lose any "
        "pedagogical accuracy. The Socratic state-routing task is a classification "
        "problem, not a generation problem.\n\n"
        "Second: ROUGE and BLEU are systematically misleading for any teaching task where "
        "the reference model was fine-tuned on the test distribution. The published winner "
        "of the KELE benchmark ranks last on pedagogical accuracy when you use a "
        "memorization-resistant metric. The unified score we propose — half stage-balanced "
        "accuracy, half LLM-judge — is one way to fix this. We don't claim it's the only "
        "way.\n\n"
        "Third: when you do measure pedagogy instead of n-gram overlap, open-weight local "
        "models like Gemma 31B reach frontier parity. Our system ties Claude Sonnet within "
        "measurement noise, runs on a consumer GPU, and costs zero API dollars per "
        "inference.\n\n"
        "The paper goes to the deeper experiments — eight-cell cross-teacher matrix, "
        "bootstrap convergence analysis showing n=400 is sufficient, the full "
        "contamination proof. We're also mid-sprint on a Stage 2 SFT teacher of our own — "
        "we shipped the pipeline yesterday in PR #79, and training is running right now "
        "on Gemma 4 31B. We had to pivot off of Qwen3.6-27B because of a Triton kernel bug "
        "on our AMD GPU — happy to talk about that in Q&A if it's interesting.\n\n"
        "Repo, HuggingFace, and the paper draft are all linked. Happy to take questions.",
    )

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = make_presentation()
    print(f"wrote {path}")
